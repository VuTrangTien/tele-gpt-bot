import logging
import os
import tempfile
from telegram import Update, InputFile
from telegram.ext import ApplicationBuilder, ContextTypes, CommandHandler, MessageHandler, filters
from telegram.constants import ChatAction
from openai import OpenAI
import openai
from docx import Document
from pptx import Presentation
import pandas as pd
import fitz  # PyMuPDF
from io import BytesIO 


# Thiết lập log
logging.basicConfig(level=logging.INFO)

# Lấy biến môi trường
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY") 

# In ra để debug
print("TELEGRAM_TOKEN =", TELEGRAM_TOKEN)
print("OPENAI_API_KEY =", OPENAI_API_KEY)

# Khởi tạo client
client = OpenAI(api_key=OPENAI_API_KEY)

# Vai trò mặc định
user_roles = {}
default_role = "Bạn là trợ lý AI hỗ trợ người dùng bằng tiếng Việt, trả lời ngắn gọn, dễ hiểu."

# Bộ nhớ tạm lưu file sửa nội dung
edited_files = {}

# ========================= HANDLERS ===============================

async def set_role(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    cmd = update.message.text

    roles = {
        "/tutor": "Bạn là giáo viên tiếng Anh, chuyên sửa lỗi bài viết, giải thích ngữ pháp và văn phong.",
        "/pm": "Bạn là chuyên gia quản lý dự án phần mềm (PMP, PMI), tư vấn hybrid, agile, waterfall, estimation...",
        "/ba": "Bạn là Business Analyst, chuyên tạo/review BRD, FSD, RFQ, RFP, Diagram, chỉnh format, ngữ pháp, phát triển ý.",
        "/sa": "Bạn là chuyên gia kiến trúc hệ thống (System Architect), thiết kế hệ thống scalable, fault-tolerant.",
        "/reset": default_role
    }

    if cmd in roles:
        user_roles[chat_id] = roles[cmd]
        await update.message.reply_text(f"✅ Đã chuyển vai trò bot sang: {cmd[1:]}")
    else:
        await update.message.reply_text("❌ Lệnh không hợp lệ.")

async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    role = user_roles.get(chat_id, default_role)
    user_message = update.message.text

    await update.message.chat.send_action(action=ChatAction.TYPING)

    try:
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": role},
                {"role": "user", "content": user_message},
            ]
        )
        reply = response.choices[0].message.content
        await update.message.reply_text(reply)

    except openai.RateLimitError:
        await update.message.reply_text("🚫 Bot đã hết quota API. Vui lòng nạp thêm để tiếp tục.")
    except Exception as e:
        await update.message.reply_text(f"⚠️ Bot lỗi: {str(e)}")

# ========================= FILE HANDLING =========================

def highlight_diff(original, edited):
    # Với văn bản đã chỉnh sửa, highlight phần khác bằng màu đỏ
    return f"\n[Trước]\n{original}\n\n[Sau]\n<font color=red>{edited}</font>"

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    role = user_roles.get(chat_id, default_role)

    file = update.message.document or update.message.effective_attachment
    file_name = file.file_name
    new_name = f"edited_{file_name}"
    file_obj = await file.get_file()
    file_bytes = await file_obj.download_as_bytearray()

    with tempfile.NamedTemporaryFile(delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    # Đọc nội dung file
    content = ""
    try:
        if file_name.endswith(".docx"):
            doc = Document(tmp_path)
            content = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        elif file_name.endswith(".pdf"):
            doc = fitz.open(tmp_path)
            content = "\n".join([page.get_text() for page in doc])
        elif file_name.endswith(".xlsx"):
            df = pd.read_excel(tmp_path)
            content = df.to_string()
        elif file_name.endswith(".pptx"):
            ppt = Presentation(tmp_path)
            content = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_name.endswith(".drawio"):
            with open(tmp_path, "r", encoding="utf-8") as f:
                content = f.read()
        else:
            await update.message.reply_text("📁 Định dạng chưa hỗ trợ.")
            return
    except Exception as e:
        await update.message.reply_text(f"❌ Lỗi đọc file: {str(e)}")
        return

    # Gửi nội dung đi để chỉnh sửa bằng OpenAI
    await update.message.chat.send_action(action=ChatAction.TYPING)
    try:
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": f"{role}\nHãy sửa lỗi, mở rộng ý và trả kết quả dạng plain text."},
                {"role": "user", "content": content},
            ]
        )
        edited_content = response.choices[0].message.content
        reply_preview = highlight_diff(content[:500], edited_content[:500])  # Preview đoạn đầu
        await update.message.reply_text(reply_preview, parse_mode="HTML")

        # Lưu để gửi lại khi người dùng yêu cầu
        edited_files[chat_id] = (file_name, edited_content)
        await update.message.reply_text("✅ Soạn thảo xong. Gõ /sendfile để nhận file đã chỉnh.")

    except Exception as e:
        await update.message.reply_text(f"⚠️ Lỗi xử lý: {str(e)}")

async def sendfile(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.effective_chat.id
    if chat_id not in edited_files:
        await update.message.reply_text("❗ Không có file nào được chỉnh gần đây.")
        return

    file_name, edited_content = edited_files[chat_id]
    out_stream = BytesIO()

    # Xuất lại đúng định dạng
    if file_name.endswith(".docx"):
        doc = Document()
        doc.add_paragraph(edited_content)
        doc.save(out_stream)
    elif file_name.endswith(".pdf"):
        pdf = fitz.open()
        page = pdf.new_page()
        page.insert_text((72, 72), edited_content)
        pdf.save(out_stream)
    elif file_name.endswith(".xlsx"):
        df = pd.DataFrame([[line] for line in edited_content.split("\n")])
        with pd.ExcelWriter(out_stream, engine="xlsxwriter") as writer:
            df.to_excel(writer, index=False, header=False)
    elif file_name.endswith(".pptx"):
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "AI Revised Content"
        slide.placeholders[1].text = edited_content[:1000]
        prs.save(out_stream)
    elif file_name.endswith(".drawio"):
        out_stream.write(edited_content.encode("utf-8"))
    else:
        await update.message.reply_text("❌ Không thể xuất định dạng này.")
        return

    out_stream.seek(0)
    await update.message.reply_document(InputFile(out_stream, filename=f"edited_{file_name}"))

# =========================== MAIN ==============================

def main():
    app = ApplicationBuilder().token(TELEGRAM_TOKEN).build()

    app.add_handler(CommandHandler(["tutor", "pm", "ba", "sa", "reset"], set_role))
    app.add_handler(CommandHandler("sendfile", sendfile))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))

    print("🤖 Bot is running...")
    app.run_polling()

if __name__ == "__main__":
    main()
